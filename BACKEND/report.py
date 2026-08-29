# builds the PDF findings report and the executive PPTX briefing off the same summary data
import os
import csv
import io
import json
from datetime import datetime
from xml.sax.saxutils import escape as _xml_escape
from reportlab.lib.pagesizes import A4
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import cm
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle,
    HRFlowable, PageBreak
)
import config

# wrap every dynamic string before it hits Paragraph — reportlab parses text as xml-ish markup
def _esc(s):
    return _xml_escape(str(s)) if s is not None else ""

RMF_COLORS = {
    "VERY_HIGH": colors.HexColor("#7B241C"),
    "HIGH":      colors.HexColor("#C0392B"),
    "MEDIUM":    colors.HexColor("#E67E22"),
    "MINOR":     colors.HexColor("#F1C40F"),
    "LOW":       colors.HexColor("#27AE60"),
    "NOT_SCORED":colors.HexColor("#95A5A6"),
}
STATUS_COLORS = {
    "GAP":                   colors.HexColor("#C0392B"),
    "PARTIAL":               colors.HexColor("#E67E22"),
    "COMPLIANT":             colors.HexColor("#27AE60"),
    "INSUFFICIENT_EVIDENCE": colors.HexColor("#7F8C8D"),
    "NOT_ASSESSED":          colors.HexColor("#BDC3C7"),
}


CSV_COLUMNS = [
    "Control ID", "Source Workbook", "Source Worksheet", "Source Row", "Source Cell",
    "Requirement", "Vendor Response", "Evidence References", "Risk Categories",
    "Risk Description", "Cause", "Consequence", "Consequence Category",
    "Existing Controls", "Control Effectiveness",
    "Initial Likelihood", "Initial Impact", "Initial Risk Rating",
    "Proposed Treatment", "Proposed Controls",
    "Residual Likelihood", "Residual Impact", "Residual Risk Rating",
    "HECVAT Compliance", "Policy Alignment", "Overall Status",
    "Exact Policy Clause", "Vendor SOC 2 Evidence", "Vendor Evidence State", "Evidence Quality",
    "Assessment Status", "Consistency Status", "Manual Review Status", "Audit Metadata",
    "Assessment Runs", "Consensus Decision",
]


def findings_to_csv(findings: list[dict]) -> str:
    buf = io.StringIO()
    w = csv.writer(buf, quoting=csv.QUOTE_ALL, lineterminator="\r\n")
    w.writerow(CSV_COLUMNS)
    for f in findings:
        w.writerow([
            f.get("control_id", ""),
            f.get("source_workbook", ""),
            f.get("source_worksheet", ""),
            f.get("source_row", ""),
            f.get("source_cell", ""),
            f.get("requirement", f.get("requirement_summary", "")),
            f.get("vendor_response", f.get("vendor_response_summary", "")),
            json.dumps(f.get("evidence_references", []), ensure_ascii=False),
            json.dumps(f.get("risk_categories", []), ensure_ascii=False),
            f.get("risk_description", ""),
            f.get("cause", ""),
            f.get("consequence", ""),
            f.get("consequence_category", ""),
            json.dumps(f.get("existing_controls", []), ensure_ascii=False),
            f.get("control_effectiveness", ""),
            f.get("initial_likelihood", 0),
            f.get("initial_impact", 0),
            config.rmf_display(f.get("initial_risk_rating", "")),
            f.get("proposed_treatment", ""),
            json.dumps(f.get("proposed_controls", []), ensure_ascii=False),
            f.get("residual_likelihood", f.get("likelihood", 0)),
            f.get("residual_impact", f.get("impact", 0)),
            config.rmf_display(f.get("residual_risk_rating", f.get("rmf_level", ""))),
            f.get("hecvat_compliance", ""),
            f.get("policy_alignment", ""),
            f.get("overall_status", ""),
            f.get("policy_clause_referenced") or "",
            json.dumps(f.get("vendor_soc2_evidence", []), ensure_ascii=False),
            f.get("vendor_evidence_state", ""),
            f.get("evidence_quality", ""),
            f.get("assessment_status", ""),
            f.get("consistency_status", ""),
            f.get("manual_review_status", ""),
            json.dumps(f.get("audit_metadata", {}), ensure_ascii=False),
            json.dumps(f.get("assessment_runs", []), ensure_ascii=False),
            json.dumps(f.get("consensus_decision", {}), ensure_ascii=False),
        ])
    # BOM so Excel opens the UTF-8 correctly instead of mangling accented vendor names
    return "\ufeff" + buf.getvalue()


def _styles():
    base = getSampleStyleSheet()
    return {
        "title": ParagraphStyle("title", parent=base["Title"],
                                fontSize=20, textColor=colors.HexColor("#1A1A2E"), spaceAfter=6),
        "h1":    ParagraphStyle("h1", parent=base["Heading1"],
                                fontSize=14, textColor=colors.HexColor("#16213E"), spaceBefore=12, spaceAfter=4),
        "h2":    ParagraphStyle("h2", parent=base["Heading2"],
                                fontSize=11, textColor=colors.HexColor("#0F3460"), spaceBefore=8, spaceAfter=3),
        "body":  ParagraphStyle("body", parent=base["Normal"],
                                fontSize=9, leading=13, spaceAfter=4),
        "small": ParagraphStyle("small", parent=base["Normal"],
                                fontSize=8, leading=11, textColor=colors.HexColor("#555555")),
        "bold":  ParagraphStyle("bold", parent=base["Normal"],
                                fontSize=9, fontName="Helvetica-Bold"),
        "table_header": ParagraphStyle(
            "table_header",
            parent=base["Normal"],
            fontSize=9,
            fontName="Helvetica-Bold",
            textColor=colors.white,
        ),
    }


def _summarise(text: str, limit: int = 400) -> str:
    t = " ".join(str(text or "").split())
    if len(t) <= limit:
        return t
    cut = t.rfind(" ", 0, limit)
    return t[:cut if cut > 0 else limit].rstrip(" ,;:.") + "\u2026"


def _exec_summary_table(summary: dict, styles: dict):
    band  = config.rmf_display(summary.get("overall_risk_band", "N/A"))
    score = summary.get("overall_rmf_score", 0)
    cov   = summary.get("coverage_pct", 0)
    data = [
        ["Metric", "Value"],
        ["Total Controls",           str(summary["total_controls"])],
        ["Assessed (policy match)",  str(summary.get("assessed_controls", 0))],
        ["Insufficient Evidence",    str(summary.get("insufficient_evidence", 0))],
        ["Inconsistent Controls",    str(summary.get("inconsistent_count", 0))],
        ["Coverage",                 f"{cov}%"],
        ["Compliant",                str(summary["status_breakdown"].get("COMPLIANT", 0))],
        ["Partial",                  str(summary["total_partial"])],
        ["Gaps",                     str(summary["total_gaps"])],
        ["Very High Risks",          str(len(summary.get("very_high_risks", [])))],
        ["High Risks",               str(len(summary.get("high_risks", [])))],
        # no leading rule here — Helvetica has no U+2500 glyph and drew black boxes
        ["RMF Score (assessed)",     str(score)],
        ["Overall Risk Band",        band],
    ]
    t = Table(data, colWidths=[10*cm, 5*cm])
    t.setStyle(TableStyle([
        ("BACKGROUND",    (0, 0), (-1, 0),  colors.HexColor("#1A1A2E")),
        ("TEXTCOLOR",     (0, 0), (-1, 0),  colors.white),
        ("FONTNAME",      (0, 0), (-1, 0),  "Helvetica-Bold"),
        ("FONTSIZE",      (0, 0), (-1, -1), 9),
        ("ROWBACKGROUNDS",(0, 1), (-1, -1), [colors.HexColor("#F8F9FA"), colors.white]),
        ("GRID",          (0, 0), (-1, -1), 0.5, colors.HexColor("#DEE2E6")),
        ("FONTNAME",      (0, 11),(-1, -1), "Helvetica-Bold"),
        ("LEFTPADDING",   (0, 0), (-1, -1), 8),
        ("RIGHTPADDING",  (0, 0), (-1, -1), 8),
        ("TOPPADDING",    (0, 0), (-1, -1), 5),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
    ]))
    return t


def _rmf_legend_table(styles: dict):
    data = [
        ["RMF Level", "Likelihood × Impact", "Action"],
        [config.rmf_display("VERY_HIGH"), "High L × High I",     "Immediate action required"],
        [config.rmf_display("HIGH"),   "High L × Med I",       "Senior management attention"],
        [config.rmf_display("MEDIUM"), "Med L × Med I",        "Management responsibility"],
        [config.rmf_display("MINOR"),  "Low L × Med I",        "Monitor and review"],
        [config.rmf_display("LOW"),    "Low L × Low I",        "Routine procedures"],
    ]
    t = Table(data, colWidths=[3*cm, 6*cm, 6*cm])
    fills = [
        colors.HexColor("#7B241C"),
        colors.HexColor("#C0392B"),
        colors.HexColor("#E67E22"),
        colors.HexColor("#F1C40F"),
        colors.HexColor("#27AE60"),
    ]
    style_cmds = [
        ("BACKGROUND",  (0, 0), (-1, 0), colors.HexColor("#16213E")),
        ("TEXTCOLOR",   (0, 0), (-1, 0), colors.white),
        ("FONTNAME",    (0, 0), (-1, 0), "Helvetica-Bold"),
        ("FONTSIZE",    (0, 0), (-1,-1), 8),
        ("GRID",        (0, 0), (-1,-1), 0.3, colors.HexColor("#DEE2E6")),
        ("TOPPADDING",  (0, 0), (-1,-1), 4),
        ("BOTTOMPADDING",(0,0), (-1,-1), 4),
        ("LEFTPADDING", (0, 0), (-1,-1), 6),
    ]
    for i, fill in enumerate(fills, 1):
        style_cmds.append(("BACKGROUND", (0, i), (0, i), fill))
        style_cmds.append(("TEXTCOLOR",  (0, i), (0, i), colors.white))
    t.setStyle(TableStyle(style_cmds))
    return t


def _section_scores_table(summary: dict, styles: dict):
    section_scores = summary.get("section_rmf_scores", {})
    if not section_scores:
        return None
    data = [["Section", "Avg RMF Score", "Band"]]
    # the cell shows the display name but colouring keys off the internal band, so keep both
    band_keys = []
    for sec, score in sorted(section_scores.items(), key=lambda x: -x[1]):
        band = config.rmf_band_from_score(score)
        band_keys.append(band)
        data.append([sec, str(round(score, 2)), config.rmf_display(band)])
    t = Table(data, colWidths=[8*cm, 4*cm, 3*cm])
    style_cmds = [
        ("BACKGROUND",    (0, 0), (-1, 0), colors.HexColor("#0F3460")),
        ("TEXTCOLOR",     (0, 0), (-1, 0), colors.white),
        ("FONTNAME",      (0, 0), (-1, 0), "Helvetica-Bold"),
        ("FONTSIZE",      (0, 0), (-1,-1), 9),
        ("ROWBACKGROUNDS",(0, 1), (-1,-1), [colors.HexColor("#F8F9FA"), colors.white]),
        ("GRID",          (0, 0), (-1,-1), 0.5, colors.HexColor("#DEE2E6")),
        ("LEFTPADDING",   (0, 0), (-1,-1), 8),
        ("TOPPADDING",    (0, 0), (-1,-1), 5),
        ("BOTTOMPADDING", (0, 0), (-1,-1), 5),
    ]
    for i, key in enumerate(band_keys, 1):
        band_col = RMF_COLORS.get(key, colors.white)
        style_cmds.append(("BACKGROUND", (2, i), (2, i), band_col))
        style_cmds.append(("TEXTCOLOR",  (2, i), (2, i),
                           colors.white if key != "MEDIUM" else colors.black))
    t.setStyle(TableStyle(style_cmds))
    return t


def _findings_table(findings: list[dict], styles: dict):
    headers = ["ID", "Section", "HECVAT", "Policy", "Overall", "RMF", "Lr", "Ir",
               "Gap / Finding", "Policy Clause", "Recommendation"]
    rows = [headers]
    for f in findings:
        rows.append([
            Paragraph(_esc(f.get("control_id", "")), styles["small"]),
            Paragraph(_esc((f.get("section") or "")[:30]), styles["small"]),
            Paragraph(_esc(f.get("hecvat_compliance") or ""), styles["small"]),
            Paragraph(_esc(f.get("policy_alignment") or ""), styles["small"]),
            Paragraph(_esc(f.get("overall_status") or ""), styles["small"]),
            Paragraph(_esc(config.rmf_display(f.get("rmf_level", ""))), styles["small"]),
            Paragraph(_esc(str(f.get("residual_likelihood", f.get("likelihood", "")))), styles["small"]),
            Paragraph(_esc(str(f.get("residual_impact", f.get("impact", "")))), styles["small"]),
            Paragraph(_esc(str(f.get("gap_description") or "")), styles["small"]),
            Paragraph(_esc(str(f.get("policy_clause_referenced") or "")), styles["small"]),
            Paragraph(_esc(str(f.get("recommendation") or "")), styles["small"]),
        ])

    col_widths = [1.34*cm, 1.87*cm, 1.12*cm, 1.49*cm, 1.34*cm, 1.34*cm,
                  0.45*cm, 0.45*cm, 3.36*cm, 2.24*cm, 3.0*cm]
    t = Table(rows, colWidths=col_widths, repeatRows=1)
    style_cmds = [
        ("BACKGROUND",    (0, 0), (-1, 0), colors.HexColor("#16213E")),
        ("TEXTCOLOR",     (0, 0), (-1, 0), colors.white),
        ("FONTNAME",      (0, 0), (-1, 0), "Helvetica-Bold"),
        ("FONTSIZE",      (0, 0), (-1,-1), 7),
        ("GRID",          (0, 0), (-1,-1), 0.3, colors.HexColor("#DEE2E6")),
        ("VALIGN",        (0, 0), (-1,-1), "TOP"),
        ("TOPPADDING",    (0, 0), (-1,-1), 3),
        ("BOTTOMPADDING", (0, 0), (-1,-1), 3),
        ("LEFTPADDING",   (0, 0), (-1,-1), 3),
    ]
    for i, f in enumerate(findings, 1):
        overall = f.get("overall_status", "")
        rmf     = f.get("rmf_level", "")
        sc = STATUS_COLORS.get(overall, colors.white)
        rc = RMF_COLORS.get(rmf, colors.white)
        style_cmds.append(("BACKGROUND", (4, i), (4, i), sc))
        style_cmds.append(("TEXTCOLOR",  (4, i), (4, i), colors.white))
        style_cmds.append(("BACKGROUND", (5, i), (5, i), rc))
        style_cmds.append(("TEXTCOLOR",  (5, i), (5, i),
                           colors.white if rmf not in ("MEDIUM", "NOT_SCORED") else colors.black))
        bg = (colors.HexColor("#FFF5F5") if overall == "GAP" else
              colors.HexColor("#FFFBF0") if overall == "PARTIAL" else colors.white)
        style_cmds.append(("BACKGROUND", (0, i), (3, i), bg))
    t.setStyle(TableStyle(style_cmds))
    return t


def _finding_detail_table(finding: dict, styles: dict):
    def joined(value):
        if isinstance(value, list):
            return "; ".join(str(item) for item in value) or "None recorded"
        return str(value if value not in (None, "") else "Unavailable")

    rows = [
        ["Field", "Value"],
        ["Source", " · ".join(filter(None, [
            str(finding.get("source_workbook") or ""),
            str(finding.get("source_worksheet") or ""),
            str(finding.get("source_cell") or ""),
        ])) or "Unavailable"],
        ["Requirement", joined(finding.get("requirement") or finding.get("requirement_summary"))],
        ["Vendor response", joined(finding.get("vendor_response") or finding.get("vendor_response_summary"))],
        ["Risk categories", joined(finding.get("risk_categories"))],
        ["Risk description", joined(finding.get("risk_description"))],
        ["Cause", joined(finding.get("cause"))],
        ["Consequence", joined(finding.get("consequence"))],
        ["Existing controls", joined(finding.get("existing_controls"))],
        ["Control effectiveness", joined(finding.get("control_effectiveness"))],
        ["Initial risk", (
            f"L{finding.get('initial_likelihood', 0)} × I{finding.get('initial_impact', 0)} · "
            f"{config.rmf_display(finding.get('initial_risk_rating'))}"
        )],
        ["Treatment", joined(finding.get("proposed_treatment"))],
        ["Proposed controls", joined(finding.get("proposed_controls"))],
        ["Residual risk", (
            f"L{finding.get('residual_likelihood', 0)} × I{finding.get('residual_impact', 0)} · "
            f"{config.rmf_display(finding.get('residual_risk_rating'))}"
        )],
        ["Policy alignment", joined(finding.get("policy_alignment"))],
        ["Exact policy clause", joined(finding.get("policy_clause_referenced"))],
        ["Vendor evidence state", joined(finding.get("vendor_evidence_state"))],
        ["Evidence quality", joined(finding.get("evidence_quality"))],
        ["Evidence references", joined([
            " · ".join(filter(None, [
                str(ref.get("filename") or ""),
                f"page {ref['page']}" if ref.get("page") else "",
                str(ref.get("cell") or ""),
                str(ref.get("chunk_id") or ""),
            ]))
            for ref in finding.get("evidence_references", [])
        ])],
        ["Assessment / consistency", (
            f"{finding.get('assessment_status', 'Unavailable')} · "
            f"{finding.get('consistency_status', 'Unavailable')} · "
            f"manual review {finding.get('manual_review_status', 'Unavailable')}"
        )],
        ["Consensus decision", joined(finding.get("consensus_decision"))],
    ]
    rendered = [[
        Paragraph(
            _esc(cell),
            styles["table_header"]
            if row_index == 0
            else (styles["bold"] if column_index == 0 else styles["small"]),
        )
        for column_index, cell in enumerate(row)
    ] for row_index, row in enumerate(rows)]
    table = Table(rendered, colWidths=[4.1 * cm, 11.9 * cm], repeatRows=1)
    table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#16213E")),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.HexColor("#F8F9FA"), colors.white]),
        ("GRID", (0, 0), (-1, -1), 0.35, colors.HexColor("#DEE2E6")),
        ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ("LEFTPADDING", (0, 0), (-1, -1), 5),
        ("RIGHTPADDING", (0, 0), (-1, -1), 5),
        ("TOPPADDING", (0, 0), (-1, -1), 4),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
    ]))
    return table


def generate_pdf(findings: list[dict], summary: dict, service_name: str, output_path: str):
    doc    = SimpleDocTemplate(output_path, pagesize=A4,
                               rightMargin=1.5*cm, leftMargin=1.5*cm,
                               topMargin=2*cm, bottomMargin=2*cm)
    styles = _styles()
    story  = []
    now_str = datetime.now().strftime("%Y-%m-%d %H:%M")

    story.append(Paragraph("IT Vendor Risk Assessment Report", styles["title"]))
    story.append(Paragraph(f"Service: {_esc(service_name)}", styles["h2"]))
    story.append(Paragraph(
        f"Generated: {now_str} | Classification: INTERNAL CONFIDENTIAL | "
        f"Framework: Murdoch University Risk Management Framework",
        styles["small"]
    ))
    story.append(HRFlowable(width="100%", thickness=1, color=colors.HexColor("#1A1A2E")))
    story.append(Spacer(1, 0.3*cm))
    story.append(Paragraph(
        "Assessment criteria: Murdoch University internal policy documents and HECVAT guidance. "
        "Session-scoped vendor documents may corroborate or contradict vendor claims but do not "
        "replace institutional policy. No unapproved external assessment framework is applied. "
        "Risk scored via Murdoch RMF likelihood × impact matrix on assessed controls only. "
        "Controls marked INSUFFICIENT_EVIDENCE reflect policy knowledge-base coverage gaps, "
        "not vendor compliance status.",
        styles["body"]
    ))
    story.append(Spacer(1, 0.3*cm))

    story.append(Paragraph("1. Executive Summary", styles["h1"]))
    story.append(_exec_summary_table(summary, styles))
    story.append(Spacer(1, 0.5*cm))

    story.append(Paragraph("2. Murdoch RMF Risk Legend", styles["h1"]))
    story.append(_rmf_legend_table(styles))
    story.append(Spacer(1, 0.5*cm))

    sec_table = _section_scores_table(summary, styles)
    if sec_table:
        story.append(Paragraph("3. Risk by Section (Assessed Controls Only)", styles["h1"]))
        story.append(sec_table)
        story.append(Spacer(1, 0.5*cm))

    very_high = summary.get("very_high_risks", [])
    if very_high:
        # no emoji here: Helvetica has no glyph for U+1F534, so ReportLab renders it as a black box
        story.append(Paragraph(f"{config.rmf_display('VERY_HIGH')} Risks \u2014 Immediate Action Required", styles["h2"]))
        for g in very_high:
            story.append(Paragraph(
                f"&#x2022; [{_esc(g['control_id'])}] {_esc(_summarise(g.get('gap_description', '')))}", styles["body"]
            ))
        story.append(Spacer(1, 0.3*cm))

    high = summary.get("high_risks", [])
    if high:
        # no emoji here either — Helvetica does not include those glyphs
        story.append(Paragraph("HIGH Risks \u2014 Senior Management Attention", styles["h2"]))
        for g in high:
            story.append(Paragraph(
                f"&#x2022; [{_esc(g['control_id'])}] {_esc(_summarise(g.get('gap_description', '')))}", styles["body"]
            ))
        story.append(Spacer(1, 0.3*cm))

    story.append(PageBreak())
    story.append(Paragraph("4. Detailed Findings", styles["h1"]))
    story.append(Spacer(1, 0.2*cm))
    story.append(_findings_table(findings, styles))
    for finding in findings:
        story.append(PageBreak())
        story.append(Paragraph(
            f"Finding {_esc(finding.get('control_id', 'Unknown'))} — "
            f"{_esc(finding.get('section', 'Unknown section'))}",
            styles["h1"],
        ))
        story.append(_finding_detail_table(finding, styles))

    doc.build(story)
    print(f"  📄 PDF: {output_path}")


# below this share of controls assessed, the sample is too thin to recommend anything
MIN_COVERAGE_FOR_VERDICT = 50.0

# how many risks the top-risks slide names; caps the slide, not the assessment
TOP_RISKS_ON_DECK = 6


# deterministic from the numbers alone so the deck's verdict never contradicts the PDF/CLI output
def _derive_recommendation(summary: dict) -> dict:
    very_high = summary.get("very_high_risks", [])
    high    = summary.get("high_risks", [])
    nv, nh  = len(very_high), len(high)
    band    = summary.get("overall_risk_band", "NOT_ASSESSED")
    score   = summary.get("overall_rmf_score", 0)
    gaps    = summary.get("total_gaps", 0)
    partial = summary.get("total_partial", 0)
    compliant = summary.get("status_breakdown", {}).get("COMPLIANT", 0)
    coverage  = summary.get("coverage_pct", 0)
    insuff    = summary.get("insufficient_evidence", 0)
    assessed  = summary.get("assessed_controls", 0)

    # no evidence is not a pass — this must never read "Approve"
    if assessed == 0 or coverage < MIN_COVERAGE_FOR_VERDICT:
        verdict = "Insufficient Evidence — Cannot Recommend"
        color   = "#7B241C"
    elif nv > 0:
        verdict = "Not Recommended — Remediate Before Onboarding"
        color   = "#7B241C"
    elif band == "HIGH" or nh >= 3:
        verdict = "Conditional Approval — High Risks Require Remediation"
        color   = "#C0392B"
    # partials are unresolved gaps in all but name
    elif gaps > 0 or partial > 0 or nh > 0:
        verdict = "Approve with Conditions"
        color   = "#E67E22"
    else:
        verdict = "Approve"
        color   = "#27AE60"

    reasoning = [
        f"Overall risk rating is {config.rmf_display(band)} (average RMF {score}) across "
        f"{assessed} assessed controls.",
    ]
    if nv or nh:
        parts = []
        if nv: parts.append(f"{nv} {config.rmf_display('VERY_HIGH')}")
        if nh: parts.append(f"{nh} {config.rmf_display('HIGH')}")
        reasoning.append(
            f"{' and '.join(parts)} severity risk"
            f"{'s' if (nv + nh) != 1 else ''} require senior attention before go-live."
        )
    else:
        reasoning.append(f"No {config.rmf_display('VERY_HIGH')} or {config.rmf_display('HIGH')} severity risks were identified.")
    reasoning.append(
        f"Control posture: {compliant} compliant, {partial} partial, {gaps} gaps "
        f"of {summary.get('total_controls', 0)} HECVAT controls."
    )
    reasoning.append(
        f"Evidence coverage is {coverage}%"
        + (f" — {insuff} controls lacked sufficient policy evidence to score."
           if insuff else " of controls could be assessed against policy.")
    )

    next_steps = []
    if nv:
        ids = ", ".join(g.get("control_id", "?") for g in very_high[:6])
        next_steps.append(f"Remediate all {config.rmf_display('VERY_HIGH')} risks ({ids}) before granting production access.")
    if nh:
        next_steps.append(
            f"Assign senior owners to each {config.rmf_display('HIGH')} risk and obtain a "
            "dated remediation plan.")
    if partial:
        next_steps.append(
            f"Close the {partial} partial control gap{'s' if partial != 1 else ''} or "
            "document compensating controls.")
    if insuff:
        next_steps.append(
            f"Request additional evidence for the {insuff} "
            f"control{'s' if insuff != 1 else ''} with insufficient coverage.")
    next_steps.append("Re-assess after remediation and formally record residual-risk acceptance.")

    return {"verdict": verdict, "color": color, "reasoning": reasoning, "next_steps": next_steps}


def generate_pptx(findings: list[dict], summary: dict, service_name: str, output_path: str):
    # this is the exec-briefing deck, not a row-per-finding dump — that's what the PDF is for
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE

    def C(h):
        return RGBColor.from_string(h.lstrip("#"))

    NAVY   = C("#141428")
    NAVY2  = C("#1A1A2E")
    PANEL  = C("#20203A")
    BLUE   = C("#0F3460")
    CRIMSON= C("#C0392B")
    WHITE  = C("#FFFFFF")
    MUTED  = C("#A9B0C0")
    LIGHT  = C("#F8F9FA")
    RMF_HEX = {
        "VERY_HIGH": "#7B241C", "HIGH": "#C0392B", "MEDIUM": "#E67E22",
        "MINOR": "#F1C40F", "LOW": "#27AE60", "NOT_SCORED": "#95A5A6",
    }

    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height
    BLANK = prs.slide_layouts[6]

    def slide(bg=NAVY):
        s = prs.slides.add_slide(BLANK)
        r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
        r.fill.solid(); r.fill.fore_color.rgb = bg
        r.line.fill.background(); r.shadow.inherit = False
        s.shapes._spTree.remove(r._element)
        s.shapes._spTree.insert(2, r._element)
        return s

    def rect(s, l, t, w, h, fill, line=None, radius=None):
        shp_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
        shp = s.shapes.add_shape(shp_type, l, t, w, h)
        if fill is None:
            shp.fill.background()
        else:
            shp.fill.solid(); shp.fill.fore_color.rgb = fill
        if line is None:
            shp.line.fill.background()
        else:
            shp.line.color.rgb = line; shp.line.width = Pt(1)
        shp.shadow.inherit = False
        return shp

    def text(s, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
        # runs can be a single paragraph's tuples or a list of paragraphs, each a list of tuples
        tb = s.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        tf.margin_left = tf.margin_right = Pt(2)
        tf.margin_top = tf.margin_bottom = Pt(2)
        paras = runs if isinstance(runs[0], list) else [runs]
        for i, para in enumerate(paras):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            if isinstance(para, tuple):
                para = [para]
            for txt, size, bold, color in para:
                r = p.add_run(); r.text = txt
                r.font.size = Pt(size); r.font.bold = bold
                r.font.color.rgb = color; r.font.name = "Calibri"
        return tb

    def accent_bar(s):
        rect(s, 0, 0, SW, Pt(6), CRIMSON)

    def header(s, kicker, title):
        accent_bar(s)
        text(s, Inches(0.7), Inches(0.42), Inches(11.9), Inches(0.35),
             (kicker.upper(), 12, True, CRIMSON))
        text(s, Inches(0.7), Inches(0.72), Inches(11.9), Inches(0.7),
             (title, 26, True, WHITE))
        rect(s, Inches(0.7), Inches(1.5), Inches(11.93), Pt(1.2), C("#33334D"))

    rec  = _derive_recommendation(summary)
    now  = datetime.now().strftime("%d %B %Y")

    s = slide(NAVY)
    rect(s, 0, 0, SW, SH, None)  # transparent full-slide shape just to keep the shape stack consistent across slides
    accent_bar(s)
    rect(s, Inches(0.9), Inches(2.05), Inches(0.14), Inches(2.6), CRIMSON)
    text(s, Inches(1.25), Inches(1.95), Inches(11), Inches(0.5),
         ("IT VENDOR RISK ASSESSMENT", 15, True, CRIMSON))
    text(s, Inches(1.25), Inches(2.45), Inches(11.2), Inches(1.5),
         (service_name, 44, True, WHITE))
    text(s, Inches(1.25), Inches(3.95), Inches(11), Inches(0.5),
         ("Executive briefing for IT management", 18, False, MUTED))
    rect(s, Inches(1.28), Inches(4.75), Inches(6.6), Pt(1), C("#33334D"))
    text(s, Inches(1.25), Inches(5.0), Inches(11), Inches(1.2),
         [[(f"{now}", 13, False, MUTED)],
          [("Murdoch University Risk Management Framework  ·  Likelihood × Impact",
            13, False, MUTED)],
          [("Classification: INTERNAL CONFIDENTIAL", 12, True, C("#7A8296"))]])

    s = slide(NAVY)
    header(s, "The bottom line", "Recommendation")
    vc = C(rec["color"])
    rect(s, Inches(0.7), Inches(1.85), Inches(11.93), Inches(1.15), vc, radius=True)
    text(s, Inches(1.05), Inches(1.85), Inches(11.3), Inches(1.15),
         (rec["verdict"], 24, True, WHITE), anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(0.72), Inches(3.35), Inches(11.9), Inches(0.4),
         ("WHY", 13, True, CRIMSON))
    text(s, Inches(0.72), Inches(3.75), Inches(11.9), Inches(3.3),
         [[("•  ", 15, True, vc), (r, 15, False, LIGHT)] for r in rec["reasoning"]],
         anchor=MSO_ANCHOR.TOP)

    s = slide(NAVY)
    header(s, "Key components", "Risk posture at a glance")
    # keep both: the raw key drives the colour lookup, the display name is what is rendered
    band_key = summary.get("overall_risk_band", "N/A")
    band = config.rmf_display(band_key)
    cards = [
        ("Total controls", str(summary.get("total_controls", 0)), MUTED),
        ("Coverage", f"{summary.get('coverage_pct', 0)}%", MUTED),
        ("Compliant", str(summary.get("status_breakdown", {}).get("COMPLIANT", 0)), C("#27AE60")),
        ("Partial", str(summary.get("total_partial", 0)), C("#E67E22")),
        ("Gaps", str(summary.get("total_gaps", 0)), C("#C0392B")),
        (f"{config.rmf_display('VERY_HIGH')} risks", str(len(summary.get("very_high_risks", []))), C("#7B241C")),
        (f"{config.rmf_display('HIGH')} risks", str(len(summary.get("high_risks", []))), C("#C0392B")),
        # `band` stays the internal key for the colour lookup; only the caption is translated
        ("Overall band", band, C(RMF_HEX.get(band_key, "#95A5A6"))),
    ]
    cw, ch, gap = Inches(2.85), Inches(1.75), Inches(0.3)
    x0, y0 = Inches(0.7), Inches(1.95)
    for i, (label, val, accent) in enumerate(cards):
        col, row = i % 4, i // 4
        l = x0 + col * (cw + gap)
        t = y0 + row * (ch + Inches(0.35))
        rect(s, l, t, cw, ch, PANEL, radius=True)
        rect(s, l, t, Inches(0.11), ch, accent)
        vsize = 40 if len(val) <= 4 else 26
        text(s, l + Inches(0.3), t + Inches(0.18), cw - Inches(0.45), Inches(1.0),
             (val, vsize, True, WHITE), anchor=MSO_ANCHOR.MIDDLE)
        text(s, l + Inches(0.3), t + ch - Inches(0.55), cw - Inches(0.45), Inches(0.45),
             (label.upper(), 12, True, MUTED))
    text(s, Inches(0.72), y0 + 2 * ch + Inches(0.55), Inches(11.9), Inches(0.5),
         ("Scored via Murdoch RMF on assessed controls only. Controls with "
          "insufficient policy evidence are tracked as coverage, not compliance.",
          11, False, C("#7A8296")))

    # bars are drawn as plain rectangles instead of a native pptx chart, because per-point chart fills make PowerPoint flag the file as needing repair
    s = slide(NAVY)
    header(s, "Key components", "Risk distribution by severity")
    rmf = summary.get("rmf_breakdown", {})
    order = ["VERY_HIGH", "HIGH", "MEDIUM", "MINOR", "LOW"]
    actions = {
        "VERY_HIGH": "Immediate action", "HIGH": "Senior management",
        "MEDIUM": "Management owns", "MINOR": "Monitor & review",
        "LOW": "Routine handling",
    }
    counts = {k: rmf.get(k, 0) for k in order}
    top = max(counts.values()) or 1
    ty = Inches(2.15)
    row_h = Inches(0.72)
    bx, bmax = Inches(3.0), Inches(6.2)
    for lvl in order:
        n = counts[lvl]
        col = C(RMF_HEX[lvl])
        text(s, Inches(0.7), ty, Inches(2.15), row_h,
             (config.rmf_display(lvl), 14, True, WHITE), anchor=MSO_ANCHOR.MIDDLE)
        rect(s, bx, ty + Inches(0.14), bmax, Inches(0.4), C("#2A2A45"), radius=True)
        fillw = max(Emu(int(bmax * (n / top))), Inches(0.06)) if n else Inches(0.06)
        rect(s, bx, ty + Inches(0.14), fillw, Inches(0.4), col, radius=True)
        text(s, bx + bmax + Inches(0.2), ty, Inches(0.9), row_h,
             (str(n), 16, True, WHITE), anchor=MSO_ANCHOR.MIDDLE)
        text(s, bx + bmax + Inches(1.15), ty, Inches(2.6), row_h,
             (actions[lvl], 12, False, MUTED), anchor=MSO_ANCHOR.MIDDLE)
        ty += row_h + Inches(0.12)
    text(s, Inches(0.7), ty + Inches(0.15), Inches(11.9), Inches(0.5),
         (f"{summary.get('assessed_controls', 0)} controls scored; bar length is relative to the "
          "largest severity band. Required action follows the Murdoch RMF.",
          11, False, C("#7A8296")))

    # one ranked slide, not one slide per risk — that produced a 117-slide deck on a real run.
    # Per-control detail lives in the PDF and CSV, which is where a reviewer works
    all_priority = summary.get("very_high_risks", []) + summary.get("high_risks", [])
    s = slide(NAVY)
    header(s, "Top risks requiring action",
           f"{len(all_priority)} High or Very High residual risk{'s' if len(all_priority) != 1 else ''}"
           if all_priority else "No High or Very High residual risks")

    if not all_priority:
        text(s, Inches(0.7), Inches(3.2), Inches(11.9), Inches(1.0),
             ("No High or Very High residual risks were identified in this assessment.",
              18, False, LIGHT))
    else:
        shown = all_priority[:TOP_RISKS_ON_DECK]
        col_x = [Inches(0.7), Inches(2.25), Inches(6.35), Inches(8.35), Inches(10.5)]
        col_w = [Inches(1.45), Inches(4.0), Inches(1.9), Inches(2.05), Inches(2.13)]
        for label, x, w in zip(("CONTROL", "RISK", "RESIDUAL", "TREATMENT", "SECTION"), col_x, col_w):
            text(s, x, Inches(1.8), w, Inches(0.3), (label, 11, True, CRIMSON))
        rect(s, Inches(0.7), Inches(2.12), Inches(11.93), Pt(1), C("#33334D"))

        # six rows plus the footnote have to fit above the 7.5in canvas edge
        row_y = Inches(2.28)
        row_h = Inches(0.72)
        for finding in shown:
            level = config.normalize_rmf_level(
                finding.get("residual_risk_rating") or finding.get("rmf_level"))
            rect(s, Inches(0.7), row_y, Inches(11.93), row_h - Inches(0.08), PANEL, radius=True)
            rect(s, Inches(0.7), row_y, Inches(0.09), row_h - Inches(0.08),
                 C(RMF_HEX.get(level, "#C0392B")))
            values = (
                str(finding.get("control_id", "Unknown")),
                _summarise(str(finding.get("risk_description")
                               or finding.get("gap_description") or "Unavailable"), 110),
                config.rmf_display(level),
                str(finding.get("proposed_treatment", "REVIEW")),
                _summarise(str(finding.get("section", "")), 40),
            )
            for value, x, w in zip(values, col_x, col_w):
                text(s, x + Inches(0.15), row_y, w, row_h - Inches(0.08),
                     (value, 11, False, LIGHT), anchor=MSO_ANCHOR.MIDDLE)
            row_y += row_h

        remaining = len(all_priority) - len(shown)
        note = (
            f"Showing the {len(shown)} highest-rated of {len(all_priority)} priority risks"
            + (f"; the remaining {remaining} are listed in full in the PDF report and CSV export."
               if remaining else ". Full cause, consequence, controls and evidence for every "
               "control are in the PDF report and CSV export.")
        )
        text(s, Inches(0.7), min(row_y + Inches(0.15), Inches(6.75)), Inches(11.9), Inches(0.45),
             (note, 11, False, C("#7A8296")))

    # the limits of the evidence belong next to the verdict, not buried behind it
    s = slide(NAVY)
    header(s, "Assurance and limitations", "What this assessment does and does not establish")
    limit_cards = [
        ("Controls assessed", f"{summary.get('assessed_controls', 0)} of "
                              f"{summary.get('total_controls', 0)}", MUTED),
        ("Evidence coverage", f"{summary.get('coverage_pct', 0)}%", MUTED),
        ("Insufficient evidence", str(summary.get("insufficient_evidence", 0)), C("#E67E22")),
        ("Awaiting manual review", str(summary.get("manual_review_count", 0)), C("#C0392B")),
    ]
    cw2, ch2, gap2 = Inches(2.85), Inches(1.5), Inches(0.3)
    for i, (label, val, accent) in enumerate(limit_cards):
        l = Inches(0.7) + i * (cw2 + gap2)
        rect(s, l, Inches(1.95), cw2, ch2, PANEL, radius=True)
        rect(s, l, Inches(1.95), Inches(0.11), ch2, accent)
        text(s, l + Inches(0.3), Inches(2.1), cw2 - Inches(0.45), Inches(0.85),
             (val, 30 if len(val) <= 6 else 20, True, WHITE), anchor=MSO_ANCHOR.MIDDLE)
        text(s, l + Inches(0.3), Inches(2.95), cw2 - Inches(0.45), Inches(0.45),
             (label.upper(), 11, True, MUTED))

    caveats = [
        f"Controls with no matching Murdoch policy are recorded as not assessed, not as compliant "
        f"({summary.get('insufficient_evidence', 0)} this run).",
        f"{summary.get('inconsistent_count', 0)} control(s) did not reach a usable model result and "
        "are excluded from the scores above until a reviewer resolves them.",
        "Vendor-supplied answers are unverified claims unless corroborated by attached evidence.",
        "Findings are model-generated and require sign-off by a qualified assessor before "
        "any onboarding decision is made.",
    ]
    text(s, Inches(0.72), Inches(3.85), Inches(11.9), Inches(0.4),
         ("READ THIS ALONGSIDE THE VERDICT", 12, True, CRIMSON))
    text(s, Inches(0.72), Inches(4.25), Inches(11.9), Inches(2.5),
         [[("•  ", 14, True, CRIMSON), (c, 14, False, LIGHT)] for c in caveats],
         anchor=MSO_ANCHOR.TOP)


    s = slide(NAVY)
    header(s, "The path forward", "Recommended next steps")
    ty = Inches(2.05)
    for i, step in enumerate(rec["next_steps"], 1):
        rect(s, Inches(0.7), ty, Inches(0.55), Inches(0.55), CRIMSON, radius=True)
        text(s, Inches(0.7), ty, Inches(0.55), Inches(0.55),
             (str(i), 18, True, WHITE), align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(s, Inches(1.5), ty - Inches(0.02), Inches(11.1), Inches(0.7),
             (step, 15, False, LIGHT), anchor=MSO_ANCHOR.MIDDLE)
        ty += Inches(0.92)
    rect(s, Inches(0.7), Inches(6.85), Inches(11.93), Pt(1), C("#33334D"))
    text(s, Inches(0.7), Inches(6.95), Inches(11.9), Inches(0.4),
         ("Prepared with the SEDONA assessment engine · Murdoch University RMF", 10, False, C("#6A7286")))

    prs.save(output_path)
    print(f"  📊 PPTX: {output_path}")


def generate_all(findings: list[dict], summary: dict,
                 service_name: str, output_dir: str = ".",
                 session_id: str | None = None):
    os.makedirs(output_dir, exist_ok=True)
    # All export formats consume the same normalized canonical dataset.
    from assess import normalize_finding
    findings = [normalize_finding(finding) for finding in findings]
    # only sanitize for the filename — the display-facing service_name stays untouched
    safe = (service_name.replace("/", "_").replace("\\", "_")
            .replace(" ", "_").lower()[:40])
    ts   = datetime.now().strftime("%Y%m%d_%H%M")
    # stamping the session id into the filename lets api.py match reports back to a session after a restart, without it we'd be stuck fuzzy-matching by vendor name
    sid = "".join(c for c in str(session_id) if c.isalnum()) if session_id else ""
    sid_part = f"{sid}_" if sid else ""
    pdf_path  = os.path.join(output_dir, f"risk_assessment_{safe}_{sid_part}{ts}.pdf")
    pptx_path = os.path.join(output_dir, f"risk_briefing_{safe}_{sid_part}{ts}.pptx")
    csv_path  = os.path.join(output_dir, f"risk_assessment_{safe}_{sid_part}{ts}.csv")
    generate_pdf(findings,  summary, service_name, pdf_path)
    generate_pptx(findings, summary, service_name, pptx_path)
    # seal the CSV with the other two, or a later code change makes the three exports disagree.
    # utf-8 not utf-8-sig: findings_to_csv already prepends the BOM, and two BOMs break the header
    with open(csv_path, "w", encoding="utf-8", newline="") as fh:
        fh.write(findings_to_csv(findings))
    print(f"  📑 CSV: {csv_path}")
    return pdf_path, pptx_path, csv_path
